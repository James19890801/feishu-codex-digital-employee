#!/usr/bin/env python3
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

def env_int(name: str, default: int, minimum: int, maximum: int) -> int:
    try:
        return max(minimum, min(maximum, int(os.environ.get(name, str(default)))))
    except (TypeError, ValueError):
        return default


MAX_CHARS = env_int("AIPRO_EXTRACT_MAX_CHARS", 40000, 1, 1000000)
MAX_PDF_PAGES = env_int("AIPRO_EXTRACT_MAX_PDF_PAGES", 60, 1, 2000)
OCR_SAMPLE_PAGES = env_int("AIPRO_PDF_OCR_SAMPLE_PAGES", 20, 1, 200)
OCR_FRONTLOAD_PAGES = env_int("AIPRO_PDF_OCR_FRONTLOAD_PAGES", 8, 0, 50)


def docx_text(path: Path) -> str:
    from docx import Document

    document = Document(path)
    chunks = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            if any(values):
                chunks.append("\t".join(values))
    return "\n".join(chunks)


def select_pdf_page_indices(
    page_count: int, sample_count: int, frontload: int = 8,
) -> list[int]:
    total = max(0, int(page_count))
    count = max(0, min(total, int(sample_count)))
    if not total or not count:
        return []
    if count >= total:
        return list(range(total))
    first_count = min(count, max(0, int(frontload)))
    selected = list(range(first_count))
    remaining = count - first_count
    if remaining <= 0:
        return selected
    start = first_count
    if remaining == 1:
        candidates = [total - 1]
    else:
        candidates = [
            round(start + index * (total - 1 - start) / (remaining - 1))
            for index in range(remaining)
        ]
    return sorted(set(selected + candidates))


def run_pdf_ocr(command: list[str]) -> list[dict]:
    try:
        result = subprocess.run(
            command, capture_output=True, text=True, timeout=300, check=True,
        )
        payload = json.loads(result.stdout)
        return payload.get("pages", []) if isinstance(payload, dict) else []
    except (OSError, subprocess.SubprocessError, json.JSONDecodeError):
        return []


def pdf_ocr(path: Path, page_count: int) -> list[dict]:
    configured = os.environ.get("AIPRO_PDF_OCR_COMMAND", "").strip()
    if configured:
        return run_pdf_ocr([configured, str(path)])
    helper = Path(__file__).resolve().parent.parent / "scripts" / "extract-pdf-ocr.swift"
    if not helper.exists():
        return []
    if os.environ.get("AIPRO_PDF_OCR_DISTRIBUTED", "").strip() != "1":
        return run_pdf_ocr(["/usr/bin/swift", str(helper), str(path)])

    from pypdf import PdfReader, PdfWriter

    selected = select_pdf_page_indices(
        page_count,
        OCR_SAMPLE_PAGES,
        frontload=OCR_FRONTLOAD_PAGES,
    )
    if not selected:
        return []
    reader = PdfReader(path)
    recognized = []
    with tempfile.TemporaryDirectory(prefix="aipro-pdf-ocr-") as temporary:
        for batch_index in range(0, len(selected), 20):
            batch = selected[batch_index:batch_index + 20]
            writer = PdfWriter()
            for page_index in batch:
                writer.add_page(reader.pages[page_index])
            sampled_path = Path(temporary) / f"sample-{batch_index // 20 + 1}.pdf"
            with sampled_path.open("wb") as output:
                writer.write(output)
            for item in run_pdf_ocr(["/usr/bin/swift", str(helper), str(sampled_path)]):
                relative_page = int(item.get("page", 0) or 0) - 1
                if 0 <= relative_page < len(batch):
                    recognized.append({
                        "page": batch[relative_page] + 1,
                        "text": str(item.get("text", "")),
                    })
    return recognized


def pdf_text(path: Path) -> tuple[str, bool]:
    from pypdf import PdfReader

    reader = PdfReader(path)
    chunks = []
    needs_ocr = False
    for index, page in enumerate(reader.pages[:MAX_PDF_PAGES], start=1):
        page_text = (page.extract_text() or "").strip()
        if len(page_text) < 20:
            needs_ocr = True
        if page_text:
            chunks.append(f"[PDF 第 {index} 页]\n{page_text}")
        if sum(len(item) for item in chunks) >= MAX_CHARS:
            break
    ocr_used = False
    if needs_ocr and sum(len(item) for item in chunks) < MAX_CHARS:
        for item in pdf_ocr(path, len(reader.pages)):
            page = int(item.get("page", 0) or 0)
            text = str(item.get("text", "")).strip()
            if page > 0 and text:
                chunks.append(f"[PDF 第 {page} 页 OCR]\n{text}")
                ocr_used = True
            if sum(len(value) for value in chunks) >= MAX_CHARS:
                break
    return "\n\n".join(chunks), ocr_used


def xlsx_text(path: Path) -> str:
    from openpyxl import load_workbook

    book = load_workbook(path, read_only=True, data_only=True)
    chunks = []
    for sheet in book.worksheets:
        chunks.append(f"[工作表：{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                chunks.append("\t".join(values).rstrip())
            if sum(len(item) for item in chunks) >= MAX_CHARS:
                return "\n".join(chunks)
    return "\n".join(chunks)


def xls_text(path: Path) -> str:
    import xlrd

    book = xlrd.open_workbook(path)
    chunks = []
    for sheet in book.sheets():
        chunks.append(f"[工作表：{sheet.name}]")
        for row_index in range(sheet.nrows):
            values = [str(value) for value in sheet.row_values(row_index)]
            if any(value.strip() for value in values):
                chunks.append("\t".join(values).rstrip())
            if sum(len(item) for item in chunks) >= MAX_CHARS:
                return "\n".join(chunks)
    return "\n".join(chunks)


def pptx_text(path: Path) -> str:
    from pptx import Presentation

    deck = Presentation(path)
    chunks = []
    for index, slide in enumerate(deck.slides, start=1):
        chunks.append(f"[幻灯片 {index}]")
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                chunks.append(text)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"[备注]\n{notes}")
        except (AttributeError, KeyError):
            pass
        if sum(len(item) for item in chunks) >= MAX_CHARS:
            break
    return "\n".join(chunks)


def legacy_office_text(path: Path) -> str:
    target = "docx" if path.suffix.lower() == ".doc" else "pptx"
    configured = os.environ.get("AIPRO_OFFICE_CONVERTER", "").strip()
    bundled = Path(sys.executable).resolve().parents[2] / "bin" / "override" / "soffice"
    converter = (
        configured
        or (str(bundled) if bundled.is_file() else "")
        or shutil.which("soffice")
        or shutil.which("libreoffice")
    )
    if not converter:
        raise ValueError(
            f"legacy {path.suffix.lower()} conversion requires LibreOffice"
        )
    with tempfile.TemporaryDirectory(prefix="aipro-office-") as temporary:
        output_dir = Path(temporary)
        profile_dir = output_dir / "profile"
        command = [
            converter,
            f"-env:UserInstallation={profile_dir.as_uri()}",
            "--headless",
            "--convert-to", target,
            "--outdir", str(output_dir),
            str(path),
        ]
        result = subprocess.run(
            command, capture_output=True, text=True, timeout=60, check=False,
        )
        converted = output_dir / f"{path.stem}.{target}"
        if result.returncode != 0 or not converted.is_file():
            detail = (result.stderr or result.stdout or "conversion failed").strip()
            raise ValueError(f"legacy Office conversion failed: {detail[:300]}")
        return docx_text(converted) if target == "docx" else pptx_text(converted)


def plain_text(path: Path) -> str:
    raw = path.read_bytes()
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "big5"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="replace")


def main() -> None:
    if len(sys.argv) != 2:
        raise SystemExit("usage: extract_file_text.py FILE")
    path = Path(sys.argv[1])
    suffix = path.suffix.lower()
    metadata = {"ocrUsed": False, "officeConverted": False}
    if suffix == ".pdf":
        text, metadata["ocrUsed"] = pdf_text(path)
    elif suffix == ".docx":
        text = docx_text(path)
    elif suffix == ".xlsx":
        text = xlsx_text(path)
    elif suffix == ".xls":
        text = xls_text(path)
    elif suffix == ".pptx":
        text = pptx_text(path)
    elif suffix in {".doc", ".ppt"}:
        text = legacy_office_text(path)
        metadata["officeConverted"] = True
    elif suffix in {".txt", ".md", ".csv", ".json", ".log", ".xml", ".html"}:
        text = plain_text(path)
    else:
        raise ValueError(f"unsupported extension: {suffix or '(none)'}")
    normalized = "\n".join(line.rstrip() for line in text.splitlines()).strip()
    print(json.dumps({
        "text": normalized[:MAX_CHARS],
        "chars": len(normalized),
        **metadata,
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
